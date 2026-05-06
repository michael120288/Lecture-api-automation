from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── colour palette ──────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_RED  = RGBColor(0xc0, 0x39, 0x2b)
ACCENT_GREEN= RGBColor(0x2e, 0x7d, 0x32)
ACCENT_BLUE = RGBColor(0x15, 0x65, 0xc0)
LIGHT_GREEN = RGBColor(0xe8, 0xf5, 0xe9)
LIGHT_BLUE  = RGBColor(0xe3, 0xf2, 0xfd)
LIGHT_RED   = RGBColor(0xff, 0xeb, 0xee)
WHITE       = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY  = RGBColor(0xf5, 0xf5, 0xf5)
MID_GRAY    = RGBColor(0x78, 0x78, 0x78)
CODE_BG     = RGBColor(0x1e, 0x1e, 0x2e)
CODE_FG     = RGBColor(0xcb, 0xd6, 0xf7)

BLANK = prs.slide_layouts[6]   # completely blank


# ── helpers ─────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK_NAVY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_code(slide, code, l, t, w, h, font_size=10):
    bg = add_rect(slide, l, t, w, h, fill=CODE_BG)
    txb = slide.shapes.add_textbox(
        Inches(l + 0.15), Inches(t + 0.12),
        Inches(w - 0.3),  Inches(h - 0.24))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code.strip().split('\n'):
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = line if line else ' '
        run.font.name = 'Courier New'
        run.font.size = Pt(font_size)
        run.font.color.rgb = CODE_FG
    return txb


def add_callout(slide, text, l, t, w, h,
                bg=LIGHT_GREEN, border=ACCENT_GREEN, size=14):
    add_rect(slide, l, t, w, h, fill=bg, line=2, line_color=border)
    txb = slide.shapes.add_textbox(
        Inches(l + 0.18), Inches(t + 0.1),
        Inches(w - 0.36), Inches(h - 0.2))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = DARK_NAVY
    return txb


def slide_header(slide, title, subtitle=None):
    """Dark navy top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill=DARK_NAVY)
    add_text(slide, title, 0.35, 0.15, 12.5, 0.75,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.78, 12.5, 0.35,
                 size=14, color=RGBColor(0xaa, 0xbb, 0xcc))


def add_bullet_list(slide, items, l, t, w,
                    size=16, color=DARK_NAVY, indent="  "):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(4))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = ("• " if not item.startswith("  ") else "  ◦ ") + item.lstrip()
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_table(slide, headers, rows, l, t, w, col_widths=None,
              header_bg=DARK_NAVY, row_bg_alt=LIGHT_GRAY):
    cols = len(headers)
    total_rows = len(rows) + 1
    row_h = 0.38
    tbl = slide.shapes.add_table(
        total_rows, cols,
        Inches(l), Inches(t),
        Inches(w), Inches(row_h * total_rows)).table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    # header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = h
        run.font.bold = True; run.font.size = Pt(13)
        run.font.color.rgb = WHITE

    # data rows
    for ri, row in enumerate(rows):
        bg = row_bg_alt if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = str(val)
            run.font.size = Pt(12); run.font.color.rgb = DARK_NAVY
    return tbl


def footer(slide, text="Lecture 01 — API Automation · codeandtest.com"):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill=DARK_NAVY)
    add_text(slide, text, 0.3, 7.17, 12.7, 0.3,
             size=10, color=RGBColor(0x99, 0xaa, 0xbb))


# ════════════════════════════════════════════════════════════════
# SLIDE 01 — Title
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK_NAVY)
add_rect(sl, 0, 3.5, 13.33, 0.06, fill=ACCENT_RED)

add_text(sl, "Lecture 01", 1, 1.2, 11, 1,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Setup & Your First API Test", 1, 2.4, 11, 0.7,
         size=28, color=RGBColor(0xaa, 0xcc, 0xff), align=PP_ALIGN.CENTER)
add_text(sl, "API Automation with Vitest + Axios + TypeScript",
         1, 3.2, 11, 0.5, size=18,
         color=RGBColor(0x88, 0x99, 0xbb), align=PP_ALIGN.CENTER)

add_rect(sl, 3.5, 4.3, 6.3, 0.7, fill=ACCENT_RED)
add_text(sl, "⏱  75–90 min  ·  codeandtest.com",
         3.5, 4.35, 6.3, 0.6, size=16,
         color=WHITE, align=PP_ALIGN.CENTER, bold=True)

add_text(sl, "🎤  Ask: who has written an automated test before? Bridge from there.",
         0.5, 6.5, 12.3, 0.6, size=11,
         color=RGBColor(0x66, 0x77, 0x88), italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 02 — Agenda
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "📋  Agenda — What We'll Cover Today")
add_table(sl,
    ["#", "Topic", "⏱ Time"],
    [
        ["1", "Why API testing + our tools",       "10 min"],
        ["2", "▶ Project setup",                   "20 min"],
        ["3", "Key concepts: validateStatus, async","15 min"],
        ["4", "▶ Write & run your first test",     "15 min"],
        ["5", "Rate limiting + common mistakes",   "10 min"],
        ["6", "Homework walkthrough",               "5 min"],
    ],
    l=0.5, t=1.3, w=12.3,
    col_widths=[0.6, 8.5, 2.0])

add_callout(sl,
    "🎯  By the end — 6 tests passing against the real production API",
    0.5, 5.9, 12.3, 0.65,
    bg=LIGHT_GREEN, border=ACCENT_GREEN, size=15)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 03 — Why API Testing
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "🤔  Why Automate API Tests?")

add_rect(sl, 0.3, 1.25, 6.0, 3.4, fill=LIGHT_RED, line=1.5, line_color=ACCENT_RED)
add_text(sl, "❌  The Problem", 0.5, 1.3, 5.5, 0.4, size=15, bold=True, color=ACCENT_RED)
add_bullet_list(sl, [
    "100 endpoints = hours of manual Postman clicking",
    "Manual checks after every code change = not realistic",
    "Easy to miss edge cases by hand",
], l=0.55, t=1.75, w=5.5, size=15)

add_rect(sl, 6.8, 1.25, 6.0, 3.4, fill=LIGHT_GREEN, line=1.5, line_color=ACCENT_GREEN)
add_text(sl, "✅  The Solution", 7.0, 1.3, 5.5, 0.4, size=15, bold=True, color=ACCENT_GREEN)
add_bullet_list(sl, [
    "Automated tests run in seconds",
    "Catch regressions after every commit",
    "Test valid AND invalid input consistently",
], l=7.05, t=1.75, w=5.5, size=15)

add_callout(sl,
    "✅ Positive test — correct input → success response\n"
    "❌ Negative test — bad input → correct error response",
    0.3, 5.0, 12.7, 0.85, bg=LIGHT_BLUE, border=ACCENT_BLUE, size=14)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 04 — Testing Pyramid
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "📐  Where API Tests Fit — The Testing Pyramid")

pyramid = [
    (4.5, 1.3, 4.3, 0.7, ACCENT_RED,   "E2E Tests — few · slow · fragile"),
    (2.5, 2.1, 8.3, 0.8, ACCENT_BLUE,  "▶  API Tests ← we are here — many · fast · realistic"),
    (0.5, 3.0,12.3, 0.7, DARK_NAVY,    "Unit Tests — most · fastest · isolated"),
]
for lx, tx, wx, hx, col, label in pyramid:
    add_rect(sl, lx, tx, wx, hx, fill=col)
    add_text(sl, label, lx+0.15, tx+0.1, wx-0.3, hx-0.15,
             size=14, bold=(col==ACCENT_BLUE), color=WHITE)

add_callout(sl,
    "💡  API tests hit a real server over real HTTP — no mocks, no browser, no flaky UI",
    0.5, 4.1, 12.3, 0.65, bg=LIGHT_BLUE, border=ACCENT_BLUE, size=14)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 05 — Tech Stack
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "🛠️  Our Tools & Why We Chose Them")
add_table(sl,
    ["Tool", "Role", "Why Not the Alternative"],
    [
        ["Vitest",      "Test runner",       "Zero-config TypeScript · 2–5× faster than Jest"],
        ["Axios",       "HTTP client",       "Works against any URL — not just localhost"],
        ["TypeScript",  "Language",          "Catches type errors before tests run"],
        ["Faker.js",    "Test data",         "Unique username every run — no database clashes"],
        ["Postman",     "Manual exploration","Try the endpoint by hand before automating"],
    ],
    l=0.5, t=1.3, w=12.3,
    col_widths=[2.0, 2.5, 7.3])

add_callout(sl,
    "💡  Supertest requires the Express app object. "
    "Axios works against production or localhost with the same code.",
    0.5, 5.9, 12.3, 0.65, bg=LIGHT_BLUE, border=ACCENT_BLUE, size=14)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 06 — DO THIS: Steps 1-3
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "▶  Project Setup — Steps 1–3")
add_rect(sl, 0, 1.1, 13.33, 0.45, fill=ACCENT_GREEN)
add_text(sl, "🖥️  Open your terminal — type these commands",
         0.3, 1.15, 12.7, 0.35, size=14, bold=True, color=WHITE)

add_code(sl, """mkdir chatty-api-tests
cd chatty-api-tests
npm init -y""",
    l=0.4, t=1.7, w=5.9, h=1.0)

add_code(sl, """npm install axios dotenv
npm install --save-dev vitest typescript @types/node @faker-js/faker""",
    l=0.4, t=2.85, w=12.5, h=0.85)

add_table(sl,
    ["Package", "Where", "Reason"],
    [
        ["axios, dotenv",                    "dependencies",    "Needed at runtime"],
        ["vitest, typescript, @faker-js/faker","devDependencies","Dev only — not deployed"],
    ],
    l=0.4, t=3.9, w=12.5, col_widths=[4.0, 3.5, 4.5])

add_callout(sl,
    "🛑  STOP — everyone types. Two commands are intentional: "
    "makes the deps vs devDeps distinction visible.",
    0.4, 5.65, 12.5, 0.65, bg=LIGHT_GREEN, border=ACCENT_GREEN, size=13)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 07 — DO THIS: package.json scripts
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "▶  Add Scripts to package.json")
add_rect(sl, 0, 1.1, 13.33, 0.45, fill=ACCENT_GREEN)
add_text(sl, "📝  Open package.json — replace the \"scripts\" section",
         0.3, 1.15, 12.7, 0.35, size=14, bold=True, color=WHITE)

add_code(sl, """{
  "scripts": {
    "test":          "vitest run",
    "test:watch":    "vitest",
    "test:ui":       "vitest --ui",
    "test:coverage": "vitest run --coverage"
  }
}""",
    l=0.4, t=1.7, w=6.0, h=2.3)

add_table(sl,
    ["Script", "When to Use"],
    [
        ["npm test",              "Run once — use in CI"],
        ["npm run test:watch",    "Watch mode while coding"],
        ["npm run test:ui",       "Visual browser UI"],
        ["npm run test:coverage", "Coverage report"],
    ],
    l=6.7, t=1.7, w=6.2, col_widths=[3.2, 3.0])

add_callout(sl,
    "🛑  STOP — everyone adds scripts. "
    "`npm test` will be the most-typed command in this entire course.",
    0.4, 5.6, 12.5, 0.65, bg=LIGHT_GREEN, border=ACCENT_GREEN, size=13)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 08 — DO THIS: tsconfig.json
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "▶  Create tsconfig.json")
add_rect(sl, 0, 1.1, 13.33, 0.45, fill=ACCENT_GREEN)
add_text(sl, "📝  Create this file in your project root",
         0.3, 1.15, 12.7, 0.35, size=14, bold=True, color=WHITE)

add_code(sl, """{
  "compilerOptions": {
    "target": "ES2022",
    "module": "ESNext",
    "moduleResolution": "bundler",
    "strict": true,
    "esModuleInterop": true,
    "skipLibCheck": true,
    "types": ["vitest/globals"]
  },
  "include": ["src/**/*", "tests/**/*"]
}""",
    l=0.4, t=1.7, w=6.5, h=3.5)

add_text(sl, "Key option:", 7.2, 1.75, 5.8, 0.4, size=15, bold=True)
add_callout(sl,
    '"types": ["vitest/globals"]\n\n'
    "Gives your IDE autocomplete for:\n"
    "describe · it · expect · beforeAll · afterAll\n\n"
    "Without it → red squiggles in the IDE\n"
    "Tests still run — this is purely for TypeScript",
    7.2, 2.2, 5.8, 2.8, bg=LIGHT_BLUE, border=ACCENT_BLUE, size=13)

add_callout(sl,
    "🛑  STOP — everyone creates tsconfig.json",
    0.4, 5.6, 12.5, 0.55, bg=LIGHT_GREEN, border=ACCENT_GREEN, size=13)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 09 — DO THIS: vitest.config.ts
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "▶  Create vitest.config.ts")
add_rect(sl, 0, 1.1, 13.33, 0.45, fill=ACCENT_GREEN)
add_text(sl, "📝  Create this file in your project root",
         0.3, 1.15, 12.7, 0.35, size=14, bold=True, color=WHITE)

add_code(sl, """import { defineConfig } from 'vitest/config';
import { config as dotenvConfig } from 'dotenv';
import { resolve } from 'path';

dotenvConfig({ path: resolve(__dirname, '.env') });

export default defineConfig({
  test: {
    globals: true,
    testTimeout: 15000,
    reporters: ['verbose'],
    fileParallelism: false,
    env: {
      BASE_URL:      process.env.BASE_URL      ?? '',
      TEST_USERNAME: process.env.TEST_USERNAME ?? '',
      TEST_PASSWORD: process.env.TEST_PASSWORD ?? '',
    },
  },
});""",
    l=0.4, t=1.7, w=7.5, h=4.5, font_size=9.5)

add_callout(sl,
    "fileParallelism: false\n"
    "Runs one test file at a time.\n"
    "Without this → 20+ auth requests fire at once → nginx rate limit → failures",
    7.3, 2.0, 5.6, 1.6, bg=LIGHT_RED, border=ACCENT_RED, size=13)

add_callout(sl,
    "dotenvConfig() → loads .env into the main process\n"
    "env:{} → forwards vars to worker threads\n"
    "Both required — workers don't inherit process.env",
    7.3, 3.8, 5.6, 1.6, bg=LIGHT_BLUE, border=ACCENT_BLUE, size=13)

add_callout(sl,
    "🛑  STOP — everyone creates vitest.config.ts",
    0.4, 6.45, 12.5, 0.55, bg=LIGHT_GREEN, border=ACCENT_GREEN, size=13)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — DO THIS: .env + config.ts
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "▶  Create .env and src/config.ts")
add_rect(sl, 0, 1.1, 13.33, 0.45, fill=ACCENT_GREEN)
add_text(sl, "📝  Create both files now",
         0.3, 1.15, 12.7, 0.35, size=14, bold=True, color=WHITE)

add_text(sl, ".env  (project root)", 0.4, 1.7, 5.8, 0.35, size=14, bold=True)
add_code(sl, "BASE_URL=https://api.codeandtest.com/api/v1",
    l=0.4, t=2.1, w=6.0, h=0.5)
add_callout(sl,
    "⚠️  Use api.codeandtest.com\n"
    "NOT codeandtest.com — that's Vercel frontend (returns 405)\n\n"
    "🔒  Never commit .env — it's already in .gitignore",
    0.4, 2.75, 6.0, 1.5, bg=LIGHT_RED, border=ACCENT_RED, size=13)

add_text(sl, "src/config.ts", 7.0, 1.7, 6.0, 0.35, size=14, bold=True)
add_code(sl, """const BASE_URL = process.env.BASE_URL;

if (!BASE_URL) {
  throw new Error(
    'Missing env var: BASE_URL'
    + ' — copy .env.example to .env'
  );
}

export const config = { BASE_URL } as const;""",
    l=7.0, t=2.1, w=6.0, h=2.8, font_size=10)

add_callout(sl,
    "The guard throws immediately with a clear message\n"
    "if .env is missing — not a silent undefined in a test",
    7.0, 5.1, 6.0, 0.9, bg=LIGHT_BLUE, border=ACCENT_BLUE, size=13)

add_callout(sl,
    "🛑  STOP — everyone creates both files",
    0.4, 6.45, 12.5, 0.55, bg=LIGHT_GREEN, border=ACCENT_GREEN, size=13)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — CHECK: Run the test
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "✅  Check Your Work — Run the Test")
add_rect(sl, 0, 1.1, 13.33, 0.45, fill=ACCENT_BLUE)
add_text(sl, "🖥️  Run this in your terminal now",
         0.3, 1.15, 12.7, 0.35, size=14, bold=True, color=WHITE)

add_code(sl, "npm test tests/lecture-01/lecture.test.ts",
    l=0.4, t=1.7, w=12.5, h=0.55)

add_code(sl, """✓ tests/lecture-01/lecture.test.ts (6)
  ✓ Style 1: async/await  > POST /signin returns 400
  ✓ Style 1: async/await  > body has message field
  ✓ Style 1: async/await  > status field is "error"
  ✓ Style 2: .then()      > POST /signin returns 400
  ✓ Style 2: .then()      > body has message field
  ✓ Style 2: .then()      > status field is "error"

Tests  6 passed (6)   Duration  ~1.5s""",
    l=0.4, t=2.4, w=8.5, h=2.8)

add_table(sl,
    ["Error", "Cause", "Fix"],
    [
        ["405 Method Not Allowed", "Wrong domain",     "Add api. prefix to BASE_URL"],
        ["Missing env var: BASE_URL","No .env file",   "cp .env.example .env"],
        ["ECONNREFUSED",           "Wrong URL format", "Check BASE_URL has https://"],
    ],
    l=9.2, t=2.4, w=3.8, col_widths=[1.8,1.0,1.0])

add_text(sl, "Troubleshooting →", 9.15, 2.25, 3.5, 0.3, size=12, bold=True, color=ACCENT_RED)

add_callout(sl,
    "🛑  Don't move on until most of the room is green — "
    "this is the first real checkpoint of the course.",
    0.4, 5.6, 12.5, 0.65, bg=LIGHT_GREEN, border=ACCENT_GREEN, size=13)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Concept: validateStatus
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "💡  The Most Important Axios Option")

add_rect(sl, 0.3, 1.25, 6.1, 4.5, fill=LIGHT_RED, line=2, line_color=ACCENT_RED)
add_text(sl, "❌  WITHOUT — test crashes", 0.5, 1.3, 5.7, 0.45,
         size=14, bold=True, color=ACCENT_RED)
add_code(sl, """// Axios THROWS on any 4xx/5xx
const res = await axios.post(url, badData);
// AxiosError: Request failed with status 400
// expect() NEVER runs
// You get a confusing stack trace
// instead of a clear assertion failure""",
    l=0.45, t=1.85, w=5.7, h=2.5, font_size=10)

add_rect(sl, 6.9, 1.25, 6.1, 4.5, fill=LIGHT_GREEN, line=2, line_color=ACCENT_GREEN)
add_text(sl, "✅  WITH — assertion runs correctly", 7.1, 1.3, 5.7, 0.45,
         size=14, bold=True, color=ACCENT_GREEN)
add_code(sl, """// You get the response — always
const res = await axios.post(url, badData, {
  validateStatus: () => true,
});

// Now this runs correctly:
expect(res.status).toBe(400); // ✓ passes""",
    l=6.95, t=1.85, w=5.7, h=2.5, font_size=10)

add_callout(sl,
    "🔑  Every single Axios call in this course uses  { validateStatus: () => true }",
    0.3, 6.0, 12.7, 0.6, bg=LIGHT_BLUE, border=ACCENT_BLUE, size=15)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Concept: async/await vs .then()
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "💡  Two Ways to Write the Same Test")

add_text(sl, "Style 1 — async/await  (recommended)", 0.4, 1.25, 12.5, 0.4,
         size=15, bold=True, color=ACCENT_GREEN)
add_code(sl, """it('returns 400', async () => {
  const res = await axios.post(url, data, { validateStatus: () => true });
  expect(res.status).toBe(400);
});""",
    l=0.4, t=1.7, w=12.5, h=1.2)

add_text(sl, "Style 2 — .then()  (must return the promise)", 0.4, 3.1, 12.5, 0.4,
         size=15, bold=True, color=ACCENT_BLUE)
add_code(sl, """it('returns 400', () => {
  return axios.post(url, data, { validateStatus: () => true })
    .then(res => expect(res.status).toBe(400));
});""",
    l=0.4, t=3.55, w=12.5, h=1.2)

add_callout(sl,
    "⚠️  Forget `return` in .then()  →  test silently passes with ZERO assertions running\n"
    "This is a false positive — the most dangerous kind of test failure",
    0.4, 5.0, 12.5, 0.85, bg=LIGHT_RED, border=ACCENT_RED, size=14)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Axios vs Supertest
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "💡  Axios vs Supertest — Key Differences")

add_table(sl,
    ["", "Axios", "Supertest"],
    [
        ["Response body",  "res.data",    "res.body"],
        ["Works against",  "Any URL",     "Local app only"],
        ["Requires",       "Just a URL",  "The Express app object"],
        ["JSON parsing",   "Automatic",   "Automatic"],
    ],
    l=0.5, t=1.4, w=12.3, col_widths=[3.0, 4.5, 4.3])

add_callout(sl,
    "🚨  Coming from Supertest? Your muscle memory will type res.body\n"
    "     It's res.data in Axios. This is the #1 confusion for Supertest veterans.",
    0.5, 4.4, 12.3, 1.0, bg=LIGHT_RED, border=ACCENT_RED, size=14)

add_callout(sl,
    "🍎 macOS / 🪟 Windows — Axios works identically on both platforms",
    0.5, 5.65, 12.3, 0.55, bg=LIGHT_BLUE, border=ACCENT_BLUE, size=14)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Rate Limiting
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "⚡  Rate Limiting on /signin")

add_text(sl, "How it works:", 0.5, 1.3, 12.0, 0.4, size=16, bold=True)
add_bullet_list(sl, [
    "nginx limits /signin to 5 requests per minute",
    "After 5 calls → server returns 429 (Too Many Requests), not 400",
    "Students hit this in homework almost immediately",
], l=0.5, t=1.75, w=12.0, size=15)

add_code(sl, """// ❌ Breaks after 5 runs in a minute:
expect(res.status).toBe(400);       // fails as 429

// ✅ Accepts 400 OR 429:
expectRejected(res.status);         // always passes""",
    l=0.5, t=3.1, w=7.5, h=1.6)

add_table(sl,
    ["Use Case", "Pattern"],
    [
        ["Testing error path (wrong password)", "expectRejected(res.status)"],
        ["Happy path signin (beforeAll setup)",  "Add x-test-secret header"],
    ],
    l=0.5, t=5.0, w=12.3, col_widths=[5.5, 6.3])
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — 3 Common Mistakes
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "🚨  3 Most Common Mistakes")

add_table(sl,
    ["#", "Mistake", "Symptom", "Fix"],
    [
        ["1", "Missing validateStatus",       "AxiosError crash",        "Add { validateStatus: () => true }"],
        ["2", "res.body instead of res.data", "undefined",               "Use res.data"],
        ["3", "Wrong BASE_URL domain",        "405 Method Not Allowed",  "Use api.codeandtest.com"],
    ],
    l=0.4, t=1.4, w=12.5, col_widths=[0.5, 3.5, 3.5, 4.5])

add_callout(sl,
    "🌐  codeandtest.com → Vercel frontend (serves HTML, returns 405 for POST)\n"
    "🌐  api.codeandtest.com → the real API  ← always use this as BASE_URL",
    0.4, 4.5, 12.5, 1.0, bg=LIGHT_RED, border=ACCENT_RED, size=14)

add_callout(sl,
    "🍎 macOS / 🪟 Windows — same mistakes, same fixes on both platforms",
    0.4, 5.75, 12.5, 0.55, bg=LIGHT_BLUE, border=ACCENT_BLUE, size=13)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — Homework
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "📝  Homework — 7 TODOs")

add_callout(sl, "Open:  tests/lecture-01/homework/starter.test.ts",
    0.4, 1.25, 12.5, 0.5, bg=DARK_NAVY, border=DARK_NAVY, size=14)

add_table(sl,
    ["TODO", "What to Practice"],
    [
        ["1",        "toContain on content-type header"],
        ["2",        "toMatchObject with expect.any()"],
        ["3",        ".not.toHaveProperty() — negative assertion"],
        ["4",        "Boundary value + expectRejected"],
        ["5 (bonus)","  .then() style — don't forget return"],
        ["6",        "toMatch(/regex/)"],
        ["7",        "toBeTypeOf('number') + toBeTruthy"],
    ],
    l=0.4, t=1.9, w=12.5, col_widths=[1.8, 10.2])

add_code(sl, "npm test tests/lecture-01/homework/starter.test.ts",
    l=0.4, t=5.75, w=8.5, h=0.5)
add_callout(sl, "🎯  Goal: 7 tests passing   Stuck? Open solution.test.ts",
    9.2, 5.75, 3.7, 0.5, bg=LIGHT_GREEN, border=ACCENT_GREEN, size=12)
footer(sl)


# ════════════════════════════════════════════════════════════════
# SLIDE 18 — What's Next
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK_NAVY)
add_rect(sl, 0, 3.3, 13.33, 0.06, fill=ACCENT_RED)

add_text(sl, "⏭️  Next Up — Lecture 02: Signin",
         0.5, 0.6, 12.3, 0.8, size=32, bold=True, color=WHITE)

add_text(sl, "What we'll add:", 0.5, 1.6, 12.0, 0.4,
         size=16, bold=True, color=RGBColor(0xaa, 0xcc, 0xff))
add_bullet_list(sl, [
    "Test the SUCCESS path of /signin",
    "Capture the session cookie from response headers",
    "Use the cookie to make authenticated requests",
], l=0.6, t=2.05, w=8.0, size=16, color=WHITE)

add_text(sl, "Before next lecture:", 0.5, 3.7, 6.0, 0.4,
         size=14, bold=True, color=RGBColor(0xaa, 0xcc, 0xff))
add_code(sl, """git add .
git commit -m "lecture-01: project setup and first test"
git checkout -b lecture-02-signin""",
    l=0.5, t=4.15, w=7.5, h=1.2)

add_callout(sl,
    "Same /signin endpoint — completely different assertions",
    8.5, 4.15, 4.5, 0.65, bg=RGBColor(0x2a, 0x2a, 0x4e),
    border=ACCENT_BLUE, size=14)

add_text(sl, "🍎 macOS / 🪟 Windows — git commands are identical on both platforms",
         0.5, 6.5, 12.0, 0.4, size=12,
         color=RGBColor(0x66, 0x77, 0x88), italic=True)


# ════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════
out = "/Users/michael120288/WebstormProjects/fullStack/theProject/documentation/slides/lecture-01.pptx"
prs.save(out)
print(f"Saved: {out}")
